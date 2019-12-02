# -*- coding: utf:8-*-
# !/usr/bin/env python
# @Date    : 2019-10-12 16:42
# @Author  : olaheihei (728037549@qq.com)


from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from test_code.config import *
from test_code.sqlop import *
import time
from pptx.util import Inches, Pt, Cm
from PIL import Image
import platform
import os
import datetime
import ast
from textwrap3 import wrap


def line_feed(text,num):
    new_text=""
    for i in wrap(text, num):
        new_text+=i
        new_text+='\n'
    return new_text

class serverPPTX(conPPTX, SqlOperate):

    def __init__(self):
        super(serverPPTX, self).__init__()  # 使用ppt周报配置类
        self.host = self.db['host']
        self.user = self.db['user']
        self.passwd = self.db['passwd']
        self.database = self.db['database']

    # 写库
    def write_content(self, group_id, content):
        self.dbcur()
        sql = "insert into week_cont(`group_id`,content)values(%s,'%s')" % (group_id, content)
        try:
            self.sqlExe(sql)
            self.sqlCom()
            self.sqlclo()
            ret = "add succ"
        except pymysql.err.IntegrityError as e:
            ret = e
        return ret

    # 更新
    def update_content(self, group_id, content):
        self.dbcur()
        sql = "UPDATE `week_cont` SET `content`=%s WHERE `group_id`=%s order by `id` desc LIMIT 1" % (content, group_id)
        try:
            self.sqlExe(sql)
            self.sqlCom()
            self.sqlclo()
            ret = "update succ"
        except pymysql.err.IntegrityError as e:
            ret = e
        return ret

    # 获取数据
    def get_week(self, group_id=0):
        # 先获得时间数组格式的日期
        threeDayAgo = (datetime.datetime.now() - datetime.timedelta(days=3))
        # 转换为时间戳
        # timeStamp = int(time.mktime(threeDayAgo.timetuple()))
        # 转换为其他字符串格式
        otherStyleTime = threeDayAgo.strftime("%Y-%m-%d %H:%M:%S")
        self.dbcur()
        if group_id == 0:
            sql = "select group_id,content from week_cont where create_time in (select max(create_time) from week_cont where create_time>'%s' group by group_id)" % otherStyleTime
        else:
            sql = "select group_id,content from week_cont where group_id=%d order by create_time desc limit 1" % int(
                group_id)
        self.sqlExe(sql)
        self.sqlCom()
        self.sqlclo()
        data = self.cur.fetchall()
        data = list(data)
        return data

    # md格式转为dict，生成列表用
    def md_to_dict(self, content):
        content = str(content).split("\n")
        d = {}
        n = 0
        for i in content:
            if n == 0:
                d['id'] = i[1:-1].split("|")
            if n > 1:
                d[str(n - 1)] = i[1:-1].split("|")
            n += 1
        return d

    # def team_building(self):


class playPPTX(conPPTX):
    # 初始化，创建ppt对象
    def __init__(self, path):
        super(playPPTX, self).__init__()  # 使用ppt周报配置类
        self.prs = Presentation(path)

    # 添加首页
    def homePage(self):
        t = time.strftime('%Y-%m-%d', time.localtime(time.time()))
        self._hp_add_time()
        return t

    # 绘制完成后，删除母版
    def _delPPTX_template(self):
        # print(len(self.base_layouts))
        for i in range(1, len(self.base_layouts)):
            self.prs.part.drop_rel(self.prs.slides._sldIdLst[1].rId)
            del self.prs.slides._sldIdLst[1]

    def getText(self):
        pass

    # 保存ppt
    def pptxSave(self, path, del_t=None):
        if not del_t is None:
            self._delPPTX_template()
        end_time = int(time.time())
        self.prs.save(path)

    def get_current_week(self):
        monday, friday = datetime.date.today(), datetime.date.today()
        one_day = datetime.timedelta(days=1)
        while monday.weekday() != 0:
            monday -= one_day
        while friday.weekday() != 4:
            friday += one_day
        monday = datetime.datetime.strftime(monday, "%Y-%m-%d")
        friday = datetime.datetime.strftime(friday, "%Y-%m-%d")
        # 返回当前的星期一和星期天的日期
        return monday, friday

    # 生成当前周的首页
    def _hp_add_time(self):
        self.prs.slides.add_slide(self.prs.slide_layouts[0])
        shapes = self.prs.slides[0].shapes  # 第一张幻灯片
        left = Inches(9.5)
        top = Inches(5.7)
        width = Inches(3)
        height = Inches(0.7)
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        content = tf.add_paragraph()
        content.text = self.get_current_week()[0] + "--" + self.get_current_week()[1]
        content.font.bold = False
        content.font.name = self.font_name
        content.font.size = Pt(16)

    # 添加目录页
    def addIndex(self):
        self.prs.slides.add_slide(self.prs.slide_layouts[1])

    # 添加表格页
    def addTable(self, work_dict, title_1):
        title = {}
        title['id'] = work_dict['id']
        work_dict.pop('id')
        size = int(len(work_dict) / 10)
        mm = len(work_dict) % 10
        cols = len(list(work_dict.values())[0]) + 1
        if size == 0:
            if mm == 0:
                raise mm
            else:
                z = 1
        else:
            if mm == 0:
                z = size
            else:
                z = size + 1

        for i in range(z):
            if i + 1 == z and mm != 0:
                work_dict = self._addTable_data(self._addTable_shapes(title, mm + 1, cols), work_dict)
                self.addTitle(title_1)
            else:
                work_dict = self._addTable_data(self._addTable_shapes(title, 11, cols), work_dict)
                self.addTitle(title_1)

    def ellipsis_str(self, str, num=16):
        if len(str) > num:
            return str[:14] + "..."
        else:
            return str

    def _addTable_shapes(self, title, rows, cols):
        left = Cm(1.0)
        top = Cm(3.0)
        width = Cm(6.0)
        height = Cm(0.8)
        shapes = self.prs.slides.add_slide(self.prs.slide_layouts[2]).shapes
        table = shapes.add_table(rows, cols, left, top, width, height).table
        if cols == 6:
            table.columns[0].width = Cm(1.2)
            table.columns[1].width = Cm(8.0)
            table.columns[2].width = Cm(4.0)
            table.columns[3].width = Cm(4.0)
            table.columns[4].width = Cm(3.0)
            table.columns[5].width = Cm(3.0)
        elif cols == 7:
            table.columns[0].width = Cm(1.2)
            table.columns[1].width = Cm(8.0)
            table.columns[2].width = Cm(4.0)
            table.columns[3].width = Cm(4.0)
            table.columns[5].width = Cm(3.0)
            table.columns[6].width = Cm(5.5)
            table.columns[7].width = Cm(3.5)
        elif cols == 8:
            table.columns[0].width = Cm(1.2)
            table.columns[1].width = Cm(8.0)
            table.columns[2].width = Cm(4.0)
            table.columns[3].width = Cm(4.0)
            table.columns[4].width = Cm(4.0)
            table.columns[5].width = Cm(3.0)
            table.columns[6].width = Cm(3.5)
            table.columns[7].width = Cm(3.5)
        else:
            raise Exception("Invalid cols")

        for r, cs in title.items():  # 绘制title
            y = 0
            # table.cell(0, y).text = r
            # print(dir(table.cell(0, y)))
            # 此处为处理表格内数据样式
            tf = table.cell(0, y).text_frame
            content = tf.paragraphs[0]
            content.text = r
            content.font.bold = True
            content.font.name = self.font_name
            content.font.size = Pt(14)

            for c in cs:
                y += 1
                tf = table.cell(0, y).text_frame
                content = tf.paragraphs[0]
                content.text = c
                content.font.bold = False
                content.font.name = self.font_name
                content.font.size = Pt(14)

        return table

    def _addTable_data(self, table, work_dict):

        a = []
        x = 1
        # cell(x,y) x 为行,y 为列 
        for r, cs in work_dict.items():  # 绘制内容
            if x > 10:
                break
            y = 0
            # table.cell(x, y).text = r
            tf = table.cell(x, y).text_frame
            content = tf.paragraphs[0]
            content.text = self.ellipsis_str(r)
            content.font.bold = False
            content.font.name = self.font_name
            content.font.size = Pt(12)
            n = 0
            for c in cs:
                if n == 0:
                    data_text = self.ellipsis_str(c, 30)
                else:
                    data_text = self.ellipsis_str(c)
                y += 1
                # table.cell(x, y).text = c
                tf = table.cell(x, y).text_frame
                content = tf.paragraphs[0]
                content.text = data_text
                content.font.bold = False
                content.font.name = self.font_name
                content.font.size = Pt(12)
                n += 1
            x += 1
            a.append(r)
        for i in a:
            work_dict.pop(i)
        return work_dict

    # 增加图片页
    def addImage(self, img_path, title):
        # Inches(1) 约等于 Image.size 110
        width, height = self._imgaeSize(img_path)
        width, height = self._calculateSize(round(width / 110, 2), round(height / 110, 2))
        blank_slide_layout = self.prs.slide_layouts[4]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        top = Inches(2.5)
        if width > height:
            left = Inches(round((13.35 - width) / 2, 2))
            width = Inches(width)
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        else:
            left = Inches(round((13.35 - width) / 2, 2))
            height = Inches(height)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
        self.addTitle(title)

    def _calculateSize(self, width, height):
        if width > 12:
            width, height = 11, round(11 / width, 2) * height

        if height > 4.5:
            width, height = round((4 / height) * width, 2), 4

        return width, height

    def _imgaeSize(self, img_path):

        return Image.open(img_path).size

    def addText(self, position, text, bold=False, i=None, size=14):
        """
        :param i: 模板的第几页,默认不添加新页
        :param position: 位置
        :param text: 文本内容
        :param bold: 是否加粗
        :return:
        """
        if i != None:
            self.prs.slides.add_slide(self.prs.slide_layouts[i])
        shapes = self.prs.slides[-1].shapes  # 第一张幻灯片
        left = Inches(position[0])
        top = Inches(position[1])
        width = Inches(position[2])
        height = Inches(position[3])
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        content = tf.paragraphs[0]
        content.text = text
        content.font.bold = bold
        content.font.name = self.font_name
        content.font.size = Pt(size)
        # content.hyperlink.address('https://www.baidu.com')

    def addTitle(self, title, i=None):
        self.addText((0.4, 0.10, 1, 1), title, True, i, size=18)

    # 数据统计
    def data_statis(self):
        self.prs.slides.add_slide(self.prs.slide_layouts[3])
        shapes = self.prs.slides[-3].shapes  # 第一张幻灯片
        left = Inches(4.5)
        top = Inches(1.7)
        width = Inches(5)
        height = Inches(5)
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        content = tf.paragraphs[0]
        content.text = "上周测试项目："
        content.font.bold = True
        content.font.name = self.font_name
        content.font.size = Pt(16)

        content = tf.paragraphs[0]
        content.text = "7个"
        content.font.bold = False
        content.font.name = self.font_name
        content.font.size = Pt(16)
        #


class playMethods():

    def imgUp(self, f):
        way = "../ppt/img/"
        basepath = os.path.dirname(__file__)  # 当前文件所在路径
        upload_path = os.path.abspath(os.path.join(basepath, way, str(int(time.time())) + '.png'))
        f.save(upload_path)

        return str(int(time.time())) + '.png'


class generatePPTX(conPPTX):
    def __init__(self, group_id=0):
        super(generatePPTX, self).__init__()  # 使用ppt周报配置类
        self.play = playPPTX('ppt/t.pptx')
        self.play.homePage()
        # self.play.addIndex()
        self.S = serverPPTX()
        self.dataAll = self.S.get_week(group_id)

    # 上周工作内容表格
    def generate_table(self, data):
        title = self.group_name + data['weekiy_name']
        content = self.S.md_to_dict(data['content'])
        self.play.addTable(content, title)

    # 团队建设
    def generate_team(self, data):
        title = self.group_name + data['weekiy_name']
        self.play.addTitle(title, 3)
        team_data = str(data['content'])[2:].split('\n##')
        if len(team_data) == 3:
            a = 0.5
            b = 1
        elif len(team_data) == 4:
            a = 0.3
            b = 0.7
        else:
            raise Exception("团队建设至少两种类型")
        top = 1.7
        for i in team_data:
            self.play.addText((4.5, top, 7, 2), i.split('\n')[0], bold=True, size=18)
            top += a
            self.play.addText((4.5, top, 7, 2), "\n".join(i.split('\n')[1:]), bold=False, size=14)
            top += b

    # 自定义文本
    def generate_diy_text(self, data):
        title = self.group_name + data['weekiy_name']
        self.play.addTitle(title, 3)
        self.play.addText((4.5, 1.7, 7, 2), line_feed(data['content'],36), bold=False, size=14)

    # 图片
    def generate_img(self, data):
        title = self.group_name + data['weekiy_name']
        self.play.addImage("ppt/img/" + data['img_name'], title)
        # num=len(data['describe'])*0.4375
        # left=round((13.35-num)/2,2)
        self.play.addText((1.5, 1.5, 7, 2), line_feed(data['describe'], 36), bold=False, size=14)

    # 问题&建议
    def generate_problem(self, data):
        self.play.addTitle(self.group_name + "问题&建议", 5)
        self.play.addText((1, 1.8, 3, 4), line_feed(data['problem'],16))
        self.play.addText((9, 1.8, 3, 4), line_feed(data['advice'],16))

    def generate(self):
        # group_id = 1
        for i in self.dataAll:
            # if i[0] != group_id:
            #     raise Exception("group_id:%d not found" % group_id)
            self.group_name = self.base_order[i[0] - 1] + "--"
            dataOne = i[1].replace('\n', '\\n')
            # 字符串转字典
            weekly = ast.literal_eval(dataOne)['weekly']

            try:
                old_table = weekly[0]['old']
                self.generate_table(old_table)
            except Exception as e:
                print(e)
                raise Exception("生成失败，上周工作内容有误")
            try:
                new_table = weekly[0]['new']
                self.generate_table(new_table)
            except Exception as e:
                print(e)
                raise Exception("生成失败，本周工作内容有误")
            try:
                steam_data = weekly[1]
                self.generate_team(steam_data)
            except Exception as e:
                print(e)
                raise Exception("生成失败，团队建设内容有误")
            try:
                problem_data = weekly[3]
                self.generate_problem(problem_data)
            except Exception as e:
                print(e)
                raise Exception("生成失败，问题&建议内容有误")

            self.play.pptxSave('t1.pptx')

    def my_sort(self, data, group_id):
        layout = self.team_layouts[self.base_order[group_id - 1]]
        L = []
        for i in layout:
            for j in data:
                if i == int(j['content_type']):
                    L.append(j)
                    break
        return L

    def generate_one(self):
        for i in self.dataAll:
            self.group_name = self.base_order[i[0] - 1] + "--"
            dataOne = i[1].replace('\n', '\\n')
            # 字符串转字典
            weekly = ast.literal_eval(dataOne)['weekly']
            other = ast.literal_eval(dataOne)['other']
            weekdata = weekly + other
            weekdata = self.my_sort(weekdata, i[0])
            for i in weekdata:
                if int(i['content_type']) == 1:
                    try:
                        old_table = i['old']
                        self.generate_table(old_table)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，上周工作内容有误")

                    try:
                        new_table = i['new']
                        self.generate_table(new_table)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，本周工作内容有误")

                elif int(i['content_type']) == 2:
                    try:
                        steam_data = i
                        self.generate_team(steam_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，团队建设内容有误")

                elif int(i['content_type']) == 3:
                    try:
                        img_data = i
                        if img_data['img_name'] != "":
                            self.generate_img(img_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，上传的图片异常")

                elif int(i['content_type']) == 4:
                    try:
                        problem_data = i
                        self.generate_problem(problem_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，问题&建议内容有误")

                elif int(i['content_type']) == 5:
                    try:
                        table_data = i
                        if table_data['content'] != "":
                            self.generate_table(table_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，自定义表格内容有误")

                elif int(i['content_type']) == 6:
                    try:
                        steam_data = i
                        if steam_data['content']!="":
                            self.generate_diy_text(steam_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，自定义文本内容有误")

                elif int(i['content_type']) == 7:
                    try:
                        img_data = i
                        if img_data['img_name'] != "":
                            self.generate_img(img_data)
                    except Exception as e:
                        print(e)
                        raise Exception("生成失败，上传的图片异常")

            today = datetime.date.today()
            today = datetime.datetime.strftime(today, "%Y-%m-%d")
            pptname = self.group_name + '周报-' + today + '.pptx'

            self.play.pptxSave('ppt/pptx/' + pptname)
            return pptname


if __name__ == '__main__':
    run = generatePPTX(11)
    run.generate_one()
